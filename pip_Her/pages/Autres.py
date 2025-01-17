#https://plotly.streamlit.app/~/+/Figure_Factory_Subplots
#https://arnaudmiribel.github.io/streamlit-extras/extras/annotated_text/
#https://docs.streamlit.io/develop/api-reference/layout/st.columns

import streamlit as st
import numpy as np
import pandas as pd
import plotly.express as px

import streamlit as st
import plotly.express as px
import pandas as pd
from fpdf import FPDF  # For PDF generation
from pptx import Presentation  # For PowerPoint generation
from docx import Document  # For Word document generation

# Dummy DataFrame
df = px.data.gapminder()

# Header
st.title("Rapport : Analyse de la diversité/inclusion")

# Sidebar for Graph Selection
graph_options = st.sidebar.multiselect(
    "Choisissez les graphiques à inclure :",
    options=["Bar Chart", "Line Chart", "Scatter Plot"],
    default=["Bar Chart"]
)

# Dictionary to store comments
comments = {}

# Generate Graphs and Comments Section
for graph_type in graph_options:
    st.subheader(graph_type)
    
    if graph_type == "Bar Chart":
        fig_bar = px.bar(df[df["year"] == 2007], x="continent", y="pop", color="continent", title="Bar Chart")
        st.plotly_chart(fig_bar, use_container_width=True)
    elif graph_type == "Line Chart":
        fig_line = px.line(df[df["country"].isin(["France", "Germany"])], x="year", y="pop", color="country", title="Line Chart")
        st.plotly_chart(fig_line, use_container_width=True)
    elif graph_type == "Scatter Plot":
        fig_scatter = px.scatter(df[df["year"] == 2007], x="gdpPercap", y="lifeExp", size="pop", color="continent", title="Scatter Plot")
        st.plotly_chart(fig_scatter, use_container_width=True)
    
    # Comment Box
    comments[graph_type] = st.text_area(f"Ajouter un commentaire pour {graph_type} :", "")

# Buttons for Export
def exporter_rapport():
    """Fonction qui transofrme la page courante dans un format choisi et donne la possibilité de la télécharger dasn ce format"""
    st.sidebar.header("Exporter le rapport")
    export_format = st.sidebar.radio("Format :", ["PDF", "Word", "PowerPoint"])

    if st.sidebar.button("Télécharger le rapport"):
        if export_format == "PDF":
            # Generate PDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.cell(200, 10, txt="Rapport Généré", ln=True, align="C")

            for graph_type, comment in comments.items():
                pdf.set_font("Arial", size=10)
                pdf.cell(0, 10, txt=f"Graphique : {graph_type}", ln=True)
                pdf.multi_cell(0, 10, txt=f" {comment}")
            
            pdf.output("rapport.pdf")
            st.sidebar.success("Le rapport PDF a été généré.")
            st.sidebar.download_button("Télécharger le PDF", data=open("rapport.pdf", "rb"), file_name="rapport.pdf")
        elif export_format == "Word":
            # Generate Word Document
            doc = Document()
            doc.add_heading("Rapport Généré", level=1)

            for graph_type, comment in comments.items():
                doc.add_heading(f"Graphique : {graph_type}", level=2)
                doc.add_paragraph(f"{comment}")
            
            doc.save("rapport.docx")
            st.sidebar.success("Le rapport Word a été généré.")
            st.sidebar.download_button("Télécharger le Word", data=open("rapport.docx", "rb"), file_name="rapport.docx")
        elif export_format == "PowerPoint":
            # Generate PowerPoint
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            title = slide.shapes.title
            title.text = "Rapport Généré"

            for graph_type, comment in comments.items():
                slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                slide.shapes.title.text = f"Graphique : {graph_type}"
                slide.placeholders[1].text = f"{comment}"
            
            presentation.save("rapport.pptx")
            st.sidebar.success("Le rapport PowerPoint a été généré.")
            st.sidebar.download_button("Télécharger le PowerPoint", data=open("rapport.pptx", "rb"), file_name="rapport.pptx")
exporter_rapport()
